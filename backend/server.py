import os
import uuid
from typing import Literal, Optional
import psycopg
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from ollama import chat
from pathlib import Path
from pypdf import PdfReader  # for PDFs [web:886]
from pptx import Presentation  # for PPTX (python-pptx)
from typing import Any, cast

DATABASE_URL = os.environ.get(
    "DATABASE_URL",
    "postgresql://postgres:postgres@localhost:5432/jorge_ai",
)

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten in prod
    allow_methods=["*"],
    allow_headers=["*"],
)


def get_conn():
    return psycopg.connect(DATABASE_URL)


class CreateConversationBody(BaseModel):
    user_id: str
    title: Optional[str] = None


class SendMessageBody(BaseModel):
    content: str
    model: str = "qwen3:4b"
    intent: Optional[
        Literal["summary", "study_plan", "practice_questions", "custom"]
    ] = None
    output_mode: Optional[Literal["quick", "full", "study_ready"]] = None


class EditConversationBody(BaseModel):
    title: str


def extract_text_from_pptx(path: str) -> str:
    pres = Presentation(path)
    parts: list[str] = []

    for idx, slide in enumerate(pres.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            sh = cast(Any, shape)  # Pylance workaround
            if not sh.has_text_frame:  # documented [web:935]
                continue

            text = sh.text_frame.text  # documented [web:935]
            if text and text.strip():
                slide_parts.append(text.strip())

        if slide_parts:
            parts.append(f"SLIDE {idx}:\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def extract_text_for_file(mime: str, path: str) -> str:
    p = Path(path)

    if mime == "application/pdf":
        reader = PdfReader(str(p))  # [web:886]
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if (
        mime
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        return extract_text_from_pptx(str(p))

    # images: not processed yet
    return ""


def build_files_context(conversation_id: str, max_chars: int = 12000) -> str:
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            SELECT filename, mime_type, storage_path
            FROM conversation_files
            WHERE conversation_id = %s
            ORDER BY id ASC
            """,
            (conversation_id,),
        )
        rows = cur.fetchall()

    blocks: list[str] = []
    for filename, mime, storage_path in rows:
        text = extract_text_for_file(mime or "", storage_path)
        if text and text.strip():
            blocks.append(f"FILE: {filename}\n{text.strip()}")

    ctx = "\n\n---\n\n".join(blocks)
    return ctx[:max_chars]


def build_system_prompt(
    intent: Literal["summary", "study_plan", "practice_questions", "custom"],
    output_mode: Literal["quick", "full", "study_ready"],
    files_text: str,
) -> str:
    mode_rules = {
        "quick": (
            "Output mode: quick.\n"
            "- Keep it compact and mobile-friendly.\n"
            "- Use short bullet points.\n"
            "- Include exactly 3 key takeaways.\n"
            "- Include exactly 3 quick practice questions."
        ),
        "full": (
            "Output mode: full.\n"
            "- Use sections in this order: Summary, Key Concepts, Examples, Common Pitfalls, Practice Questions.\n"
            "- Keep paragraphs short and scannable.\n"
            "- Include 5 practice questions with concise answers."
        ),
        "study_ready": (
            "Output mode: study_ready.\n"
            "- Use sections in this order: Summary, 7-Day Study Plan, Flashcards, Quiz, Revision Checklist.\n"
            "- Keep it practical and exam-focused.\n"
            "- Flashcards should be Q/A pairs.\n"
            "- Quiz should include answers and 1-line explanations."
        ),
    }

    intent_rules = {
        "summary": "Intent focus: prioritize a clean summary and major takeaways before anything else.",
        "study_plan": "Intent focus: prioritize a concrete study plan with realistic pacing and milestones.",
        "practice_questions": "Intent focus: prioritize high-quality exam-style questions and answers.",
        "custom": "Intent focus: answer the user question directly and then add useful study follow-ups.",
    }

    grounding_rules = (
        "You are a study assistant.\n"
        "Grounding requirements:\n"
        "- Base your answer on the uploaded files whenever possible.\n"
        "- If information is missing from files, explicitly say what is missing.\n"
        "- Do not invent slide content.\n"
        "- When useful, reference slide numbers (e.g., SLIDE 3).\n"
        "- End with 1-2 concise clarifying questions if user intent is ambiguous.\n"
    )

    files_block = (
        files_text.strip() or "No attached files were found for this conversation."
    )

    return (
        grounding_rules
        + "\n"
        + mode_rules[output_mode]
        + "\n"
        + intent_rules[intent]
        + "\n\nCourse file context:\n"
        + files_block
    )


@app.get("/conversations")
async def list_conversation(user_id: str):
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            SELECT id, title, created_at, updated_at
            FROM conversations
            WHERE user_id = %s
            ORDER BY updated_at DESC
            """,
            (user_id,),
        )
        rows = cur.fetchall()

    return {
        "conversations": [
            {
                "id": str(r[0]),
                "title": r[1],
                "created_at": r[2].isoformat(),
                "updated_at": r[3].isoformat(),
            }
            for r in rows
        ]
    }


@app.post("/conversations")
async def create_conversation(body: CreateConversationBody):
    convo_id = str(uuid.uuid4())

    try:
        with get_conn() as conn, conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO conversations (id, user_id, title, created_at, updated_at)
                VALUES (%s, %s, %s, now(), now())
                RETURNING id, title, created_at, updated_at
                """,
                (convo_id, body.user_id, body.title),
            )
            row = cur.fetchone()
            conn.commit()
    except psycopg.errors.ForeignKeyViolation:
        raise HTTPException(
            status_code=400, detail="user_id does not exist in users table"
        )

    if row is None:
        raise HTTPException(status_code=500, detail="Failed to create conversation")

    return {
        "conversation": {
            "id": str(row[0]),
            "title": row[1],
            "created_at": row[2].isoformat(),
            "updated_at": row[3].isoformat(),
        }
    }


@app.get("/conversations/{conversation_id}/messages")
async def fetch_messages(
    conversation_id: str, limit: int = 50, before_id: Optional[int] = None
):
    with get_conn() as conn, conn.cursor() as cur:
        if before_id is None:
            cur.execute(
                """
                SELECT id, role, content, created_at
                FROM messages
                WHERE conversation_id = %s
                ORDER BY id DESC
                LIMIT %s
                """,
                (conversation_id, limit),
            )
        else:
            cur.execute(
                """
                SELECT id, role, content, created_at
                FROM messages
                WHERE conversation_id = %s AND id < %s
                ORDER BY id DESC
                LIMIT %s
                """,
                (conversation_id, before_id, limit),
            )
        rows = cur.fetchall()

    rows = list(reversed(rows))  # ascending for UI
    return {
        "messages": [
            {"id": r[0], "role": r[1], "content": r[2], "created_at": r[3].isoformat()}
            for r in rows
        ]
    }


@app.post("/conversations/{conversation_id}/messages")
async def send_message(conversation_id: str, body: SendMessageBody):
    user_text = body.content.strip()
    if not user_text:
        raise HTTPException(status_code=400, detail="content is required")

    # 1) Insert user message + load context
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            INSERT INTO messages (conversation_id, role, content, created_at)
            VALUES (%s, 'user', %s, now())
            RETURNING id, created_at
            """,
            (conversation_id, user_text),
        )
        user_row = cur.fetchone()
        if user_row is None:
            raise HTTPException(status_code=500, detail="Failed to insert user message")

        cur.execute(
            """
            SELECT role, content
            FROM messages
            WHERE conversation_id = %s
            ORDER BY id DESC
            LIMIT 20
            """,
            (conversation_id,),
        )
        ctx = list(reversed(cur.fetchall()))
        conn.commit()
        files_text = build_files_context(
            conversation_id=conversation_id, max_chars=12000
        )

        intent = body.intent or "custom"
        output_mode = body.output_mode or "full"
        system_prompt = build_system_prompt(
            intent=intent, output_mode=output_mode, files_text=files_text
        )

        ollama_messages = [
            {"role": role, "content": content} for (role, content) in ctx
        ]
        ollama_messages = [
            {"role": "system", "content": system_prompt}
        ] + ollama_messages

        resp = chat(model=body.model, messages=ollama_messages, stream=False)

    assistant_text = resp.message.content

    # 3) Insert assistant + bump updated_at
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            INSERT INTO messages (conversation_id, role, content, created_at)
            VALUES (%s, 'assistant', %s, now())
            RETURNING id, created_at
            """,
            (conversation_id, assistant_text),
        )
        asst_row = cur.fetchone()
        if asst_row is None:
            raise HTTPException(
                status_code=500, detail="Failed to insert assistant message"
            )

        cur.execute(
            "UPDATE conversations SET updated_at = now() WHERE id = %s",
            (conversation_id,),
        )
        conn.commit()

    return {
        "user_message": {
            "id": user_row[0],
            "role": "user",
            "content": user_text,
            "created_at": user_row[1].isoformat(),
        },
        "assistant_message": {
            "id": asst_row[0],
            "role": "assistant",
            "content": assistant_text,
            "created_at": asst_row[1].isoformat(),
        },
    }


@app.patch("/conversations/{conversation_id}")
async def edit_conversation(conversation_id: str, body: EditConversationBody):
    title = body.title.strip()
    if not title:
        raise HTTPException(status_code=400, detail="title is required")

    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            UPDATE conversations
            SET title = %s,
                updated_at = now()
            WHERE id = %s
            RETURNING id, title, created_at, updated_at
            """,
            (title, conversation_id),
        )
        row = cur.fetchone()
        conn.commit()

    if row is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    return {
        "conversation": {
            "id": str(row[0]),
            "title": row[1],
            "created_at": row[2].isoformat(),
            "updated_at": row[3].isoformat(),
        }
    }


@app.delete("/conversations/{conversation_id}")
async def delete_conversation(conversation_id: str):

    if not conversation_id:
        raise HTTPException(status_code=400, detail="No conversation found")

    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            DELETE FROM conversations
            WHERE id = %s
            RETURNING id, title, created_at, updated_at
            """,
            (conversation_id,),
        )
        row = cur.fetchone()
        conn.commit()

    if row is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    return {
        "conversation": {
            "id": str(row[0]),
            "title": row[1],
            "created_at": row[2].isoformat(),
            "updated_at": row[3].isoformat(),
        }
    }


# --------- FILE TREATMENT -----------
UPLOAD_DIR = Path(os.environ.get("UPLOAD_DIR", "./uploads"))
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


@app.post("/conversations/{conversation_id}/files")
async def upload_conversation_file(conversation_id: str, file: UploadFile = File(...)):
    if not file.filename:
        raise HTTPException(status_code=400, detail="filename is required")

    mime = file.content_type or "application/octet-stream"

    # guarda no disco
    safe_ext = Path(file.filename).suffix.lower()
    stored_name = f"{uuid.uuid4().hex}{safe_ext}"
    dst = UPLOAD_DIR / stored_name

    size = 0
    with dst.open("wb") as f:
        while True:
            chunk = await file.read(1024 * 1024)
            if not chunk:
                break
            size += len(chunk)
            f.write(chunk)

    # regista na DB
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            INSERT INTO conversation_files (conversation_id, filename, mime_type, size_bytes, storage_path)
            VALUES (%s, %s, %s, %s, %s)
            RETURNING id, created_at
            """,
            (conversation_id, file.filename, mime, size, str(dst)),
        )
        row = cur.fetchone()
        conn.commit()

    if row is None:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return {
        "file": {
            "id": row[0],
            "conversation_id": conversation_id,
            "filename": file.filename,
            "mime_type": mime,
            "size_bytes": size,
            "created_at": row[1].isoformat(),
        }
    }


@app.get("/conversations/{conversation_id}/files")
async def list_conversation_files(conversation_id: str):
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            SELECT id, filename, mime_type, size_bytes, created_at
            FROM conversation_files
            WHERE conversation_id = %s
            ORDER BY id DESC
            """,
            (conversation_id,),
        )
        rows = cur.fetchall()

    return {
        "files": [
            {
                "id": r[0],
                "filename": r[1],
                "mime_type": r[2],
                "size_bytes": r[3],
                "created_at": r[4].isoformat(),
            }
            for r in rows
        ]
    }


@app.delete("/conversations/{conversation_id}/files/{file_id}")
async def delete_conversation_file(conversation_id: str, file_id: int):
    with get_conn() as conn, conn.cursor() as cur:
        cur.execute(
            """
            SELECT storage_path
            FROM conversation_files
            WHERE id = %s AND conversation_id = %s
            """,
            (file_id, conversation_id),
        )
        row = cur.fetchone()

        if row is None:
            raise HTTPException(status_code=404, detail="File not found")

        storage_path = row[0]

        cur.execute(
            "DELETE FROM conversation_files WHERE id = %s AND conversation_id = %s",
            (file_id, conversation_id),
        )
        conn.commit()

    # delete from disk (best effort)
    try:
        Path(storage_path).unlink(missing_ok=True)
    except Exception:
        pass

    return {"deleted": True, "file_id": file_id}
